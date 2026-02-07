"""
Generate Executive Summary PowerPoint for ASPR Photo Repository.
Uses ASPR + Leidos branding with professional slide design.

Run:  python scripts/generate_exec_summary_pptx.py
Requires: pip install python-pptx
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "ASPR_Photo_Repository_Executive_Summary.pptx"

# ── Brand Colors ──────────────────────────────────────────────────────
BLUE_DARK     = RGBColor(0x06, 0x2E, 0x61)
BLUE_PRIMARY  = RGBColor(0x15, 0x51, 0x97)
BLUE_MEDIUM   = RGBColor(0x24, 0x77, 0xBD)
GOLD          = RGBColor(0xAA, 0x64, 0x04)
RED           = RGBColor(0x99, 0x00, 0x00)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY     = RGBColor(0x32, 0x32, 0x32)
GOLD_LIGHT    = RGBColor(0xFB, 0xD0, 0x98)

# ── Logo Paths ────────────────────────────────────────────────────────
ASPR_LOGO = ROOT / "public" / "aspr-logo-blue.png"
HHS_LOGO = ROOT / "public" / "hhs_longlogo_white.png"
LEIDOS_LOGO = Path(
    r"C:\Users\ravinder.mathaudhu\OneDrive - HHS Office of the Secretary"
    r"\Documents\Projects\New folder\Leidos-Logo-Suite\Leidos-Logo-Suite"
    r"\02-Digital\03-Raster-PNG\Leidos-logo-horz-full-rgb-@2x.png"
)
OCIO_LOGO = Path(
    r"C:\Users\ravinder.mathaudhu\GitHub\ocio-techhhs-web-appservice"
    r"\public\ocio-eagle-logo.png"
)


# ── Helpers ───────────────────────────────────────────────────────────

def add_dark_bg(slide, color=BLUE_DARK):
    """Fill the entire slide with a solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.06), color=GOLD):
    """Add a thin horizontal accent bar across the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), top, Inches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, text="HHS/ASPR — For Official Use Only | Leidos"):
    """Add a footer text at the bottom of the slide."""
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12.333)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


def add_title_text(slide, text, left, top, width, height,
                   font_size=Pt(36), color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, title, bullets, title_color=WHITE, bullet_color=WHITE):
    """Add a title + bullet point content to a slide."""
    # Title
    add_title_text(slide, title,
                   Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
                   font_size=Pt(32), color=title_color, bold=True)

    # Accent bar under title
    add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)

    # Bullet content
    left = Inches(1.0)
    top = Inches(1.8)
    width = Inches(11)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = bullet_color
        p.space_after = Pt(12)
        p.level = 0

    return tf


def add_table_slide(slide, title, headers, rows, col_widths=None):
    """Add a title + table to a slide."""
    # Title
    add_title_text(slide, title,
                   Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
                   font_size=Pt(32), color=WHITE, bold=True)
    add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.5)
    height = Inches(0.4) * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(Inches(11.5) * w / total)

    # Header row
    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_PRIMARY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True

    # Data rows
    for ri, row_data in enumerate(rows):
        for ci, text in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(text)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x7A)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x08, 0x35, 0x6E)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = WHITE

    return table


# ══════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # Blank layout


# ── SLIDE 1: TITLE SLIDE ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

# Top accent bar
add_accent_bar(slide, top=Inches(0), height=Inches(0.08), color=GOLD)

# Logos row (top right area)
logo_top = Inches(0.4)
if ASPR_LOGO.exists():
    slide.shapes.add_picture(str(ASPR_LOGO), Inches(0.8), logo_top, height=Inches(1.0))
if LEIDOS_LOGO.exists():
    slide.shapes.add_picture(str(LEIDOS_LOGO), Inches(10.5), logo_top, height=Inches(0.7))

# Main title area
add_title_text(slide, "Executive Summary",
               Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
               font_size=Pt(48), color=WHITE, bold=True)

add_title_text(slide, "ASPR Photo Repository Application",
               Inches(0.8), Inches(3.2), Inches(11), Inches(0.7),
               font_size=Pt(28), color=GOLD_LIGHT, bold=False)

# Divider line
add_accent_bar(slide, top=Inches(4.1), height=Inches(0.04), color=GOLD)

# Metadata
metadata_lines = [
    "U.S. Department of Health and Human Services",
    "Administration for Strategic Preparedness and Response (ASPR)",
    "",
    "Prepared by: HHS ASPR / Leidos",
    "Date: February 7, 2026",
    "Classification: For Official Use Only (FOUO)",
]
left = Inches(0.8)
top = Inches(4.5)
width = Inches(11)
height = Inches(2.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(metadata_lines):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC) if line else WHITE
    if "Department" in line or "ASPR" in line:
        p.font.color.rgb = WHITE
        p.font.size = Pt(18)

add_footer(slide)


# ── SLIDE 2: PURPOSE & MISSION ───────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_bullet_slide(slide, "Purpose & Mission", [
    "Enable ASPR field teams to securely capture, upload, and manage disaster-related "
    "photographs during incident response operations",
    "Provide rapid photo documentation capability deployable within hours of incident activation",
    "Replace ad-hoc photo collection methods (email, shared drives, USB) with a purpose-built, "
    "secure web application",
    "Support incident accountability and documentation requirements with geotagged, "
    "timestamped photographic evidence",
    "Operate within the HHS/ASPR security boundary with full OWASP and NIST compliance",
])

add_footer(slide)


# ── SLIDE 3: KEY CAPABILITIES ────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_title_text(slide, "Key Capabilities",
               Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
               font_size=Pt(32), color=WHITE, bold=True)
add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)

# Two-column layout
capabilities_left = [
    ("Secure Authentication", "PIN-based access with bcrypt hashing,\nJWT tokens, and rate limiting"),
    ("Photo Upload Wizard", "Multi-step guided upload with\nreal-time progress tracking"),
    ("Geospatial Metadata", "GPS coordinates via device location\nor ZIP code lookup"),
]
capabilities_right = [
    ("Photo Gallery", "Review, download, and manage\nuploaded photos per session"),
    ("Admin Dashboard", "Web-based PIN creation and\nsession management"),
    ("CDN-Ready Images", "HMAC-signed URLs with\n7-day CDN cache headers"),
]

for col_idx, caps in enumerate([capabilities_left, capabilities_right]):
    x_offset = Inches(0.8) if col_idx == 0 else Inches(7.0)
    for i, (title, desc) in enumerate(caps):
        y_offset = Inches(1.9) + Inches(1.6) * i

        # Title
        add_title_text(slide, title,
                       x_offset, y_offset, Inches(5.5), Inches(0.4),
                       font_size=Pt(20), color=GOLD_LIGHT, bold=True)
        # Description
        add_title_text(slide, desc,
                       x_offset, y_offset + Inches(0.45), Inches(5.5), Inches(0.9),
                       font_size=Pt(15), color=WHITE, bold=False)

add_footer(slide)


# ── SLIDE 4: ARCHITECTURE OVERVIEW ───────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_table_slide(slide, "Architecture Overview",
    ["Component", "Technology", "Purpose"],
    [
        ["Web Framework", "Next.js 16.1.6 (React 19)", "Full-stack application"],
        ["Database", "Azure SQL Server", "Session & photo metadata"],
        ["Blob Storage", "Azure Blob Storage", "Photo file storage"],
        ["Authentication", "JWT + bcrypt + HMAC-SHA256", "Multi-layer security"],
        ["Image Processing", "Sharp", "Thumbnail generation"],
        ["Hosting", "Azure App Service (Linux)", "Node.js 22 runtime"],
        ["CI/CD", "GitHub Actions + OIDC", "Automated deployment"],
    ],
    col_widths=[25, 35, 40],
)

add_footer(slide)


# ── SLIDE 5: SECURITY POSTURE ────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_bullet_slide(slide, "Security Posture", [
    "FIPS 199 MODERATE categorization — appropriate for operational incident photography",
    "OWASP Top 10 (2021) fully addressed — injection prevention, access control, "
    "cryptographic protections",
    "NIST SP 800-63B compliant PIN generation using CSPRNG with bcrypt storage (10 salt rounds)",
    "Comprehensive rate limiting — 5 PIN attempts/min with 15-min lockout; "
    "3 admin attempts with 30-min lockout",
    "Hardened HTTP security headers — HSTS, CSP, X-Frame-Options, Permissions-Policy on all routes",
    "Zero plaintext credential storage — PINs stored as bcrypt hashes, "
    "admin tokens compared with timing-safe algorithm",
    "Signed image URLs (HMAC-SHA256) eliminate JWT exposure in query strings",
])

add_footer(slide)


# ── SLIDE 6: TIMELINE & MILESTONES ───────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_table_slide(slide, "Timeline & Milestones",
    ["Phase", "Timeline", "Status", "Deliverables"],
    [
        ["Requirements & Design", "Jan 2026", "Complete", "SRS, SDD, Security Plan"],
        ["Core Development", "Jan–Feb 2026", "Complete", "Auth, Upload, Gallery, Admin"],
        ["Security Hardening", "Feb 2026", "Complete", "bcrypt, rate limiting, signed URLs"],
        ["UI/UX Polish", "Feb 2026", "Complete", "Glassmorphic design, Lenis scroll"],
        ["Documentation", "Feb 2026", "Complete", "Full 7-document package"],
        ["UAT & Review", "Feb–Mar 2026", "Pending", "Stakeholder testing & feedback"],
        ["ATO & Production", "Mar 2026", "Planned", "Security review, go-live"],
    ],
    col_widths=[25, 15, 15, 45],
)

add_footer(slide)


# ── SLIDE 7: RISK ASSESSMENT ─────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_table_slide(slide, "Risk Assessment",
    ["Risk", "Likelihood", "Impact", "Mitigation"],
    [
        ["PIN brute force", "Low", "Medium", "Rate limiting + lockout + bcrypt"],
        ["Data loss", "Low", "High", "Azure automatic backups + blob soft delete"],
        ["Network unavailability", "Medium", "Medium", "Standalone deployment; offline upload planned"],
        ["Credential exposure", "Low", "High", "bcrypt hashes; Key Vault; no plaintext storage"],
        ["Scale limitations", "Medium", "Low", "In-memory rate limit → Redis migration path"],
        ["VNet access constraints", "Low", "Low", "Kudu-based migration scripts documented"],
    ],
    col_widths=[25, 15, 15, 45],
)

add_footer(slide)


# ── SLIDE 8: RECOMMENDATION & APPROVAL ───────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_title_text(slide, "Recommendation & Approval",
               Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
               font_size=Pt(32), color=WHITE, bold=True)
add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)

# Recommendation text
add_title_text(slide,
    "The ASPR Photo Repository application is recommended for deployment to production. "
    "The system meets all functional requirements, adheres to NIST and OWASP security standards, "
    "and has been designed for the operational needs of ASPR incident response teams.",
    Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
    font_size=Pt(18), color=WHITE, bold=False)

# Approval table
table_shape = slide.shapes.add_table(5, 4, Inches(0.8), Inches(3.4), Inches(11.5), Inches(2.5))
table = table_shape.table

headers = ["Role", "Name", "Signature", "Date"]
col_widths_pct = [30, 25, 25, 20]
total = sum(col_widths_pct)
for i, w in enumerate(col_widths_pct):
    table.columns[i].width = int(Inches(11.5) * w / total)

for i, hdr in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE_PRIMARY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True

roles = [
    "Federal Project Sponsor",
    "Information System Security Officer (ISSO)",
    "Authorizing Official (AO)",
    "Technical Lead",
]
for ri, role in enumerate(roles):
    cell = table.cell(ri + 1, 0)
    cell.text = role
    bg = RGBColor(0x0A, 0x3D, 0x7A) if ri % 2 == 0 else RGBColor(0x08, 0x35, 0x6E)
    for ci in range(4):
        c = table.cell(ri + 1, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = bg
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

add_footer(slide)


# ── SLIDE 9: CONTACT & NEXT STEPS ────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)

add_title_text(slide, "Next Steps",
               Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
               font_size=Pt(32), color=WHITE, bold=True)
add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)

next_steps = [
    "1.  Complete User Acceptance Testing (UAT) with ASPR field team representatives",
    "2.  Conduct security review and obtain Authority to Operate (ATO)",
    "3.  Configure production Azure resources (SQL, Blob Storage, Key Vault)",
    "4.  Deploy to production environment via GitHub Actions CI/CD pipeline",
    "5.  Distribute admin credentials and train operations staff on PIN management",
    "6.  Conduct field pilot during next incident activation or training exercise",
    "7.  Gather feedback and iterate on v1.1 enhancements (map view, batch download)",
]

left = Inches(1.0)
top = Inches(1.9)
width = Inches(11)
height = Inches(4.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

for i, step in enumerate(next_steps):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

# Bottom logos
if ASPR_LOGO.exists():
    slide.shapes.add_picture(str(ASPR_LOGO), Inches(0.8), Inches(6.2), height=Inches(0.7))
if LEIDOS_LOGO.exists():
    slide.shapes.add_picture(str(LEIDOS_LOGO), Inches(10.5), Inches(6.3), height=Inches(0.5))

add_footer(slide)


# ══════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
size_kb = OUT.stat().st_size / 1024
print(f"\nExecutive Summary PPTX generated: {OUT}")
print(f"Size: {size_kb:.1f} KB")
print(f"Slides: {len(prs.slides)}")
